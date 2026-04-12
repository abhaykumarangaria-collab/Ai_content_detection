from flask import Flask, request, render_template, send_file
from pptx import Presentation
from io import BytesIO
from fpdf import FPDF

# Initialize the second app
app = Flask(__name__)

# --- ROUTES ---

@app.route("/")
@app.route("/converter")
def converter(): 
    # This will load your converter.html without touching your other app
    return render_template("converter.html")

# --- CONVERTER LOGIC ---
@app.route("/api/convert", methods=["POST"])
def convert_file():
    conversion_type = request.form.get("type")
    file = request.files.get("file")

    if not file or file.filename == "":
        return "No file uploaded", 400

    try:
        # PPT to Text
        if conversion_type == "ppt-text":
            prs = Presentation(file)
            text_content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text.strip() + "\n"
                text_content += "\n--- Slide Break ---\n\n"
            
            output = BytesIO()
            output.write(text_content.encode("utf-8"))
            output.seek(0)
            return send_file(output, as_attachment=True, download_name="converted.txt", mimetype="text/plain")

        # Text to PPT
        elif conversion_type == "text-ppt":
            text_content = file.read().decode("utf-8")
            prs = Presentation()
            slides_text = text_content.split("\n\n") 
            for slide_text in slides_text:
                if slide_text.strip():
                    slide = prs.slides.add_slide(prs.slide_layouts[1]) 
                    slide.shapes.title.text = "Generated Slide"
                    slide.shapes.placeholders[1].text = slide_text.strip()
            
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            return send_file(output, as_attachment=True, download_name="converted.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        # PPT to PDF
        elif conversion_type == "ppt-pdf":
            prs = Presentation(file)
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        clean_text = shape.text.encode('latin-1', 'replace').decode('latin-1')
                        pdf.multi_cell(0, 10, clean_text)
                pdf.ln(10) 
            
            output = BytesIO()
            output.write(pdf.output(dest='S')) 
            output.seek(0)
            return send_file(output, as_attachment=True, download_name="converted.pdf", mimetype="application/pdf")

    except Exception as e:
        return f"Error during conversion: {str(e)}", 500

    return "Invalid conversion type", 400

if __name__ == "__main__": 
    # We run this on PORT 5001 so it doesn't conflict with your AI Detector on 5000
    app.run(debug=True, port=5001)