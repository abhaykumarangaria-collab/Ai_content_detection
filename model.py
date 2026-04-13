from flask import Flask, request, render_template, jsonify
from pptx import Presentation
import joblib
import re
import nltk
import xgboost as xgb
import os
import base64
import pythoncom  # Required for Windows COM threading
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize
from werkzeug.utils import secure_filename

# Optional PDF Converter setup (Requires Windows + MS Office)
try:
    import comtypes.client
    COMTYPES_AVAILABLE = True
except ImportError:
    COMTYPES_AVAILABLE = False

# 1. INITIALIZE APP
app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# 2. LOAD NLTK DATA
nltk.download("stopwords", quiet=True)
nltk.download("punkt", quiet=True)
stop_words = set(stopwords.words("english"))

# 3. LOAD MODEL & VECTORIZER
try:
    model = joblib.load("xgb_ai_model.pkl")
    vectorizer = joblib.load("vectorizer.pkl")
    print("✅ System Ready: Model & Vectorizer Loaded")
except Exception as e:
    print(f"❌ Critical Error: {e}")
    exit()

# 4. HELPER FUNCTIONS
def preprocess(text):
    text = str(text).lower()
    text = re.sub(r"[^a-zA-Z\s]", "", text)
    words = [w for w in text.split() if w not in stop_words]
    return " ".join(words)

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    full_text, slide_data = "", []
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text += shape.text + " "
        
        content = slide_text.strip()
        if content:
            slide_data.append({'slide': i + 1, 'text': content})
            full_text += content + "\n"
    return full_text.strip(), slide_data

def highlight_ai_sentences(text):
    sentences = sent_tokenize(text.replace('\n', '. ')) 
    highlighted = []
    
    for s in sentences:
        original = s.strip()
        word_count = len(original.split())
        if word_count < 6:
            highlighted.append(f"<p style='color:#94a3b8; margin: 8px 0;'>⚪ {original}</p>")
            continue
            
        clean = preprocess(original)
        vec = vectorizer.transform([clean])
        prob = 0.0 if vec.nnz == 0 else float(model.predict(xgb.DMatrix(vec))[0])
        if 0.66 < prob < 0.68: prob = 0.0
        
        color = "#ff4d4d" if prob > 0.70 else ("#ffcc00" if prob > 0.40 else "#00ffcc")
        icon = "🔴" if prob > 0.70 else ("🟡" if prob > 0.40 else "🟢")
        score_text = f"({round(prob*100,1)}%)" if prob > 0 else "(Neutral)"
        highlighted.append(f"<p style='color:{color}; margin: 8px 0;'>{icon} {original} {score_text}</p>")
        
    return "".join(highlighted)

def convert_ppt_to_pdf(input_path, output_path):
    if not COMTYPES_AVAILABLE:
        raise Exception("comtypes library missing.")
    
    # Initialize COM for this specific thread
    pythoncom.CoInitialize()
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        abs_in = os.path.abspath(input_path)
        abs_out = os.path.abspath(output_path)
        deck = powerpoint.Presentations.Open(abs_in, WithWindow=False)
        deck.SaveAs(abs_out, 32) # 32 = PDF
        deck.Close()
        powerpoint.Quit()
    finally:
        pythoncom.CoUninitialize()

# --- ROUTES ---

@app.route("/")
def home(): 
    return render_template("index.html")

@app.route("/converter")
def converter(): 
    return render_template("converter.html")

@app.route("/api/convert", methods=["POST"])
def api_convert():
    file = request.files.get("file")
    conversion_type = request.form.get("type")
    if not file or file.filename == "": return jsonify({"error": "No file"}), 400
    
    path = os.path.join(app.config["UPLOAD_FOLDER"], secure_filename(file.filename))
    file.save(path)
    
    try:
        if conversion_type == "ppt-text":
            text, _ = extract_text_from_ppt(path)
            clean_full = preprocess(text)
            vec = vectorizer.transform([clean_full])
            prob = 0.0 if vec.nnz == 0 else float(model.predict(xgb.DMatrix(vec))[0])
            b64 = base64.b64encode(text.encode("utf-8")).decode("utf-8")
            return jsonify({"status": "success", "ai_probability": round(prob * 100, 2), "human_probability": round((1-prob)*100,2), "file_b64": b64, "mime_type": "text/plain", "extension": ".txt", "scan_complete": True})

        elif conversion_type == "text-ppt":
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                lines = [line.strip() for line in f.readlines() if line.strip()]
            prs = Presentation()
            for i in range(0, len(lines), 5):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Slide {i//5 + 1}"
                tf = slide.shapes.placeholders[1].text_frame
                for j, line in enumerate(lines[i:i+5]):
                    if j == 0: tf.text = line
                    else: tf.add_paragraph().text = line
            
            out = os.path.join(app.config["UPLOAD_FOLDER"], "gen.pptx")
            prs.save(out)
            with open(out, "rb") as f: b64 = base64.b64encode(f.read()).decode("utf-8")
            
            try: os.remove(out)
            except Exception: pass
            
            return jsonify({"status": "success", "file_b64": b64, "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "extension": ".pptx", "scan_complete": False})

        elif conversion_type == "ppt-pdf":
            out = os.path.join(app.config["UPLOAD_FOLDER"], "conv.pdf")
            convert_ppt_to_pdf(path, out)
            with open(out, "rb") as f: b64 = base64.b64encode(f.read()).decode("utf-8")
            
            try: os.remove(out)
            except Exception: pass
            
            return jsonify({"status": "success", "file_b64": b64, "mime_type": "application/pdf", "extension": ".pdf", "scan_complete": False})

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        # BULLETPROOF CLEANUP
        try:
            if os.path.exists(path): 
                os.remove(path)
        except Exception:
            pass 

@app.route("/predict_ppt", methods=["POST"])
def predict_ppt():
    file = request.files.get("file")
    if not file: return "No file"
    path = os.path.join(app.config["UPLOAD_FOLDER"], secure_filename(file.filename))
    file.save(path)
    try:
        text, slide_data = extract_text_from_ppt(path)
        clean = preprocess(text)
        vec = vectorizer.transform([clean])
        prob = 0.0 if vec.nnz == 0 else float(model.predict(xgb.DMatrix(vec))[0])
        slide_scores = []
        for s in slide_data:
            c_s = preprocess(s['text'])
            s_v = vectorizer.transform([c_s])
            s_p = float(model.predict(xgb.DMatrix(s_v))[0]) if s_v.nnz > 0 and len(s['text'].split()) > 5 else 0.0
            slide_scores.append({"slide": s['slide'], "score": round(s_p * 100, 1)})
        return render_template("result.html", result="AI Detected" if prob > 0.5 else "Human Authored", probability=round(prob*100,2), highlighted_text=highlight_ai_sentences(text), slide_scores=slide_scores)
    finally:
        # BULLETPROOF CLEANUP
        try:
            if os.path.exists(path): 
                os.remove(path)
        except Exception:
            pass 

if __name__ == "__main__": 
    app.run(debug=True)