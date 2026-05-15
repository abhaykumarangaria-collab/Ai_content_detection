import os
import requests
import traceback
from dotenv import load_dotenv
from flask import Flask, request, render_template, jsonify, flash, redirect, url_for
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import joblib
import re
import nltk
import xgboost as xgb
import base64
import heapq
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize

# ==========================================
# LOAD ENV VARIABLES
# ==========================================
load_dotenv()

# ==========================================
# IMPORT RAG ENGINE
# ==========================================
try:
    from rag_engine import SlideRAGEngine
    print("✅ RAG Engine Loaded")
except Exception as e:
    print("❌ Failed to import rag_engine.py")
    print(e)

# ==========================================
# WINDOWS PDF SUPPORT
# ==========================================
try:
    import pythoncom
    import comtypes.client
    COMTYPES_AVAILABLE = True
except ImportError:
    COMTYPES_AVAILABLE = False
    print("⚠️ PDF conversion disabled")

# ==========================================
# FLASK CONFIG
# ==========================================
app = Flask(__name__)

app.config['SECRET_KEY'] = os.environ.get(
    'SECRET_KEY',
    'fallback-secret-key'
)

db_url = os.environ.get(
    'DATABASE_URL',
    'sqlite:///users.db'
)

if db_url.startswith("postgres://"):
    db_url = db_url.replace(
        "postgres://",
        "postgresql://",
        1
    )

app.config['SQLALCHEMY_DATABASE_URI'] = db_url
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

UPLOAD_FOLDER = "uploads"

app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# ==========================================
# DATABASE
# ==========================================
db = SQLAlchemy(app)

login_manager = LoginManager(app)
login_manager.login_view = 'login'

# ==========================================
# USER MODEL
# ==========================================
class User(db.Model, UserMixin):

    id = db.Column(
        db.Integer,
        primary_key=True
    )

    username = db.Column(
        db.String(150),
        unique=True,
        nullable=False
    )

    password = db.Column(
        db.String(256),
        nullable=False
    )

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

with app.app_context():
    db.create_all()

# ==========================================
# NLTK
# ==========================================
nltk.download("stopwords", quiet=True)
nltk.download("punkt", quiet=True)

stop_words = set(stopwords.words("english"))

# ==========================================
# LOAD ML MODELS
# ==========================================
print("📦 Loading ML models...")

try:
    model_xgb = joblib.load("xgb_ai_model.pkl")
    model_sgd = joblib.load("sgd_ai_model.pkl")
    vectorizer = joblib.load("vectorizer.pkl")

    print("✅ ML models loaded")

except Exception as e:
    print("❌ ML model load failed")
    print(e)

# ==========================================
# HELPER FUNCTIONS
# ==========================================
def preprocess(text):

    text = str(text).lower()

    text = re.sub(
        r"[^a-zA-Z\s]",
        "",
        text
    )

    words = [
        w for w in text.split()
        if w not in stop_words
    ]

    return " ".join(words)

def extract_text_from_ppt(file_path):

    prs = Presentation(file_path)

    full_text = ""

    slide_data = []

    for i, slide in enumerate(prs.slides):

        slide_text = ""

        for shape in slide.shapes:

            if hasattr(shape, "text") and shape.text:

                text_chunk = shape.text.strip()

                if len(text_chunk.split()) >= 5:
                    slide_text += text_chunk + " "

        content = slide_text.strip()

        if content:

            slide_data.append({
                "slide": i + 1,
                "text": content
            })

            full_text += content + "\n"

    return full_text.strip(), slide_data

def get_ensemble_score(text):

    try:
        clean = preprocess(text)

        vec = vectorizer.transform([clean])

        if vec.nnz == 0:
            return 0.0

        prob_xgb = float(
            model_xgb.predict(
                xgb.DMatrix(vec)
            )[0]
        )

        prob_sgd = float(
            model_sgd.predict_proba(vec)[0][1]
        )

        final = (prob_xgb * 0.6) + (prob_sgd * 0.4)

        return final

    except Exception as e:
        print("Score Error:", e)
        return 0.5

def analyze_tone(text):

    text_lower = str(text).lower()

    tones = {
        "Academic 🎓": [
            "research",
            "analysis",
            "methodology",
            "theory"
        ],

        "Professional 💼": [
            "strategy",
            "growth",
            "metrics"
        ],

        "Casual ☕": [
            "cool",
            "awesome",
            "fun"
        ]
    }

    scores = {
        tone: 0
        for tone in tones
    }

    words = word_tokenize(text_lower)

    for word in words:

        for tone, vocab in tones.items():

            if word in vocab:
                scores[tone] += 1

    dominant = max(scores, key=scores.get)

    if scores[dominant] == 0:
        return "Neutral 📊"

    return dominant

def generate_summary(text, num_sentences=3):

    sentences = sent_tokenize(
        text.replace("\n", ". ")
    )

    if len(sentences) <= num_sentences:
        return sentences

    words = word_tokenize(text.lower())

    word_freq = {}

    for word in words:

        if word not in stop_words and word.isalnum():

            word_freq[word] = (
                word_freq.get(word, 0) + 1
            )

    if not word_freq:
        return []

    max_freq = max(word_freq.values())

    for word in word_freq:
        word_freq[word] /= max_freq

    sentence_scores = {}

    for sent in sentences:

        for word in word_tokenize(sent.lower()):

            if word in word_freq:

                sentence_scores[sent] = (
                    sentence_scores.get(sent, 0)
                    + word_freq[word]
                )

    summary = heapq.nlargest(
        num_sentences,
        sentence_scores,
        key=sentence_scores.get
    )

    return summary

# ==========================================
# AUTH ROUTES
# ==========================================
@app.route('/register', methods=['GET', 'POST'])
def register():

    if request.method == 'POST':

        username = request.form.get('username')

        password = request.form.get('password')

        existing = User.query.filter_by(
            username=username
        ).first()

        if existing:
            flash("Username already exists")
            return redirect(url_for('register'))

        hashed = generate_password_hash(password)

        user = User(
            username=username,
            password=hashed
        )

        db.session.add(user)
        db.session.commit()

        login_user(user)

        return redirect(url_for('home'))

    return render_template("register.html")

@app.route('/login', methods=['GET', 'POST'])
def login():

    if request.method == 'POST':

        username = request.form.get('username')

        password = request.form.get('password')

        user = User.query.filter_by(
            username=username
        ).first()

        if user and check_password_hash(
            user.password,
            password
        ):

            login_user(user)

            return redirect(url_for('home'))

        flash("Invalid credentials")

    return render_template("login.html")

@app.route('/logout')
@login_required
def logout():

    logout_user()

    return redirect(url_for('login'))

# ==========================================
# PAGE ROUTES
# ==========================================
@app.route("/")
@app.route("/detector")
def home():
    return render_template("detector.html")

@app.route("/about")
def about():
    return render_template("about.html")

@app.route("/contact")
def contact():
    return render_template("contact.html")

# ==========================================
# CHAT API
# ==========================================
@app.route("/api/chat", methods=["POST"])
def api_chat():

    try:

        print("\n====================")
        print("CHAT API CALLED")
        print("====================")

        data = request.get_json(force=True)

        print("REQUEST DATA:")
        print(data)

        question = data.get("question")

        slides = data.get("slides")

        print("QUESTION:")
        print(question)

        print("SLIDES TYPE:")
        print(type(slides))

        if not question:
            return jsonify({
                "error": "Question missing"
            }), 400

        if not slides:
            return jsonify({
                "error": "Slides missing"
            }), 400

        print("Creating RAG Engine...")

        rag = SlideRAGEngine()

        print("RAG Engine Created")

        print("Ingesting Slides...")

        success = rag.ingest_slides(slides)

        print("INGEST RESULT:")
        print(success)

        if not success:

            return jsonify({
                "error": "Failed to process slides"
            }), 500

        print("Generating Answer...")

        result = rag.generate_answer(question)

        print("AI RESULT:")
        print(result)

        if not isinstance(result, dict):

            return jsonify({
                "error": "Invalid AI response format"
            }), 500

        if "error" in result:

            return jsonify({
                "error": result["error"]
            }), 500

        return jsonify({
            "answer": result.get(
                "answer",
                "No answer generated"
            )
        })

    except Exception as e:

        print("\n❌ CHAT ROUTE CRASHED")
        print(traceback.format_exc())

        return jsonify({
            "error": str(e)
        }), 500

# ==========================================
# PPT DETECTION
# ==========================================
@app.route("/predict_ppt", methods=["POST"])
def predict_ppt():

    file = request.files.get("file")

    if not file:
        return "No file uploaded"

    path = os.path.join(
        app.config["UPLOAD_FOLDER"],
        secure_filename(file.filename)
    )

    file.save(path)

    try:

        text, slide_data = extract_text_from_ppt(path)

        final_prob = get_ensemble_score(text)

        tone = analyze_tone(text)

        summary = generate_summary(text)

        return render_template(
            "result.html",
            result=(
                "AI Detected"
                if final_prob > 0.5
                else "Human Authored"
            ),
            probability=round(final_prob * 100, 2),
            tone=tone,
            summary=summary,
            raw_slides=slide_data
        )

    except Exception as e:

        print(traceback.format_exc())

        return f"Error: {str(e)}"

    finally:

        try:
            if os.path.exists(path):
                os.remove(path)
        except:
            pass

# ==========================================
# RUN APP
# ==========================================
if __name__ == "__main__":

    port = int(
        os.environ.get("PORT", 5000)
    )

    print(f"🚀 Running on port {port}")

    app.run(
        host="0.0.0.0",
        port=port,
        debug=True
    )